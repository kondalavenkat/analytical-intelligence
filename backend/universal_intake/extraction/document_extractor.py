import io
import pandas as pd
from .service import BaseExtractor
from ..models.standard_document import StandardDocument, FileType, ExtractionStrategy

MAX_PDF_PAGES = 100   # safety cap — very large PDFs handled page-by-page


class DocumentExtractor(BaseExtractor):
    """
    Handles: PDF (digital + scanned fallback), DOCX, PPTX, TXT, MD, RTF
    Extracts text and tables with explicit per-file-type error handling.
    """

    def extract(self, data: bytes, filename: str) -> StandardDocument:
        ext = filename.lower().rsplit(".", 1)[-1]
        doc = StandardDocument(
            file_type_ext=ext,
            technical_file_type=FileType.DOCUMENT,
            extraction_strategy=ExtractionStrategy.DOCUMENT,
            parser_used="DocumentExtractor",
        )

        if not data:
            doc.ok = False
            doc.flag_reason = "File is empty (0 bytes)."
            return doc

        try:
            if ext == "pdf":
                return self._parse_pdf(data, doc)
            elif ext in ("docx", "doc"):
                return self._parse_docx(data, doc, ext)
            elif ext in ("pptx", "ppt"):
                return self._parse_pptx(data, doc, ext)
            elif ext in ("txt", "md", "rtf"):
                return self._parse_txt(data, doc)
            else:
                doc.ok = False
                doc.flag_reason = f"Unsupported document type: .{ext}"
                return doc
        except Exception as e:
            import traceback
            traceback.print_exc()
            doc.ok = False
            doc.flag_reason = f"Could not process file: {e}"
            return doc

    # ── Digital PDF ───────────────────────────────────────────────────────────
    def _parse_pdf(self, data: bytes, doc: StandardDocument) -> StandardDocument:
        try:
            import pdfplumber
        except ImportError:
            doc.ok = False
            doc.flag_reason = "pdfplumber not installed. Run: pip install pdfplumber"
            return doc

        all_text:   list[str]  = []
        all_tables: list[list] = []

        try:
            with pdfplumber.open(io.BytesIO(data)) as pdf:
                if len(pdf.pages) == 0:
                    doc.ok = False
                    doc.flag_reason = "PDF has no pages (may be corrupt or password-protected)."
                    return doc

                doc.page_count = len(pdf.pages)
                pages_to_read  = pdf.pages[:MAX_PDF_PAGES]

                for page in pages_to_read:
                    try:
                        # Tables first — they are higher-quality structured data
                        tables = page.extract_tables() or []
                        for tbl in tables:
                            if tbl:  # guard None rows
                                all_tables.extend([r for r in tbl if r is not None])
                        # Raw text as fallback / supplement
                        text = page.extract_text() or ""
                        if text.strip():
                            all_text.append(text.strip())
                    except Exception as page_err:
                        print(f"[DocumentExtractor] PDF page error (skipping): {page_err}")
                        continue

        except Exception as open_err:
            # pdfplumber fails on password-protected or corrupt PDFs
            err_msg = str(open_err).lower()
            if "password" in err_msg or "encrypted" in err_msg:
                doc.ok = False
                doc.flag_reason = "PDF is password-protected. Please provide an unlocked version."
            else:
                doc.ok = False
                doc.flag_reason = f"Cannot open PDF: {open_err}"
            return doc

        doc.raw_text = "\n".join(all_text)

        if all_tables:
            # Use the first row as header — but guard empty or None-header rows
            header_row = all_tables[0]
            if not header_row or all(h is None or str(h).strip() == "" for h in header_row):
                # Unnamed headers — generate generic ones
                header_row = [f"col_{i+1}" for i in range(len(header_row or []))]
            headers = [str(h).strip() if h else f"col_{i}" for i, h in enumerate(header_row)]
            rows    = [[str(c) if c is not None else "" for c in row]
                       for row in all_tables[1:] if row]
            try:
                df = pd.DataFrame(rows, columns=headers)
                df = df.dropna(axis=1, how="all")
                doc.df        = df
                doc.row_count = len(df)
                doc.col_count = len(df.columns)
                doc.columns   = df.columns.tolist()
                doc.preview   = df.head(5).fillna("").astype(str).values.tolist()
                doc.confidence = 0.88
            except Exception as df_err:
                print(f"[DocumentExtractor] PDF table → DataFrame error: {df_err}")
                all_tables = []  # fall through to text path
        
        if not all_tables:
            if not doc.raw_text.strip():
                # Completely empty — could be a scanned PDF
                doc.ok = False
                doc.flag_reason = (
                    "No text or tables found in PDF. "
                    "If this is a scanned document, upload as an image (PNG/JPG) instead."
                )
                return doc
            lines = [l.strip() for l in doc.raw_text.split("\n") if l.strip()]
            df    = pd.DataFrame({"line_number": range(1, len(lines) + 1), "content": lines})
            doc.df        = df
            doc.row_count = len(df)
            doc.col_count = 2
            doc.columns   = ["line_number", "content"]
            doc.preview   = df.head(5).fillna("").astype(str).values.tolist()
            doc.confidence = 0.80

        print(f"[DocumentExtractor] PDF: {doc.page_count} pages, "
              f"{len(all_tables)} table rows, text_len={len(doc.raw_text)}")
        return doc

    # ── DOCX ──────────────────────────────────────────────────────────────────
    def _parse_docx(self, data: bytes, doc: StandardDocument, ext: str) -> StandardDocument:
        # Old binary .doc format is NOT supported by python-docx
        if ext == "doc":
            doc.ok = False
            doc.flag_reason = (
                "Old Word binary format (.doc) is not supported. "
                "Please save the file as .docx (File → Save As → Word Document) and re-upload."
            )
            return doc

        try:
            from docx import Document as DocxDocument
        except ImportError:
            doc.ok = False
            doc.flag_reason = "python-docx not installed. Run: pip install python-docx"
            return doc

        try:
            word_doc = DocxDocument(io.BytesIO(data))
        except Exception as open_err:
            doc.ok = False
            doc.flag_reason = f"Cannot open DOCX file (may be corrupt or password-protected): {open_err}"
            return doc

        all_text:  list[str]  = []
        table_rows: list[list] = []

        # Extract embedded tables first
        try:
            for table in word_doc.tables:
                for row in table.rows:
                    try:
                        table_rows.append([cell.text.strip() for cell in row.cells])
                    except Exception:
                        pass
        except Exception as tbl_err:
            print(f"[DocumentExtractor] DOCX table error (non-fatal): {tbl_err}")

        # Extract paragraph text
        try:
            for para in word_doc.paragraphs:
                if para.text.strip():
                    all_text.append(para.text.strip())
        except Exception as para_err:
            print(f"[DocumentExtractor] DOCX paragraph error (non-fatal): {para_err}")

        if not all_text and not table_rows:
            doc.ok = False
            doc.flag_reason = "DOCX file appears to be empty or contains only images."
            return doc

        doc.raw_text = "\n".join(all_text)

        if table_rows:
            header_row = table_rows[0]
            headers    = [str(h).strip() or f"col_{i}" for i, h in enumerate(header_row)]
            rows       = table_rows[1:]
            try:
                df = pd.DataFrame(rows, columns=headers)
                doc.df        = df
                doc.row_count = len(df)
                doc.col_count = len(df.columns)
                doc.columns   = df.columns.tolist()
                doc.preview   = df.head(5).fillna("").astype(str).values.tolist()
                doc.confidence = 0.87
            except Exception:
                table_rows = []  # fall through

        if not table_rows:
            lines = [l for l in all_text if l]
            df    = pd.DataFrame({"line_number": range(1, len(lines) + 1), "content": lines})
            doc.df        = df
            doc.row_count = len(df)
            doc.col_count = 2
            doc.columns   = ["line_number", "content"]
            doc.preview   = df.head(5).fillna("").astype(str).values.tolist()
            doc.confidence = 0.80

        print(f"[DocumentExtractor] DOCX: {len(all_text)} paragraphs, "
              f"{len(table_rows)} table rows")
        return doc

    # ── PPTX ──────────────────────────────────────────────────────────────────
    def _parse_pptx(self, data: bytes, doc: StandardDocument, ext: str) -> StandardDocument:
        # Old binary .ppt format not supported
        if ext == "ppt":
            doc.ok = False
            doc.flag_reason = (
                "Old PowerPoint binary format (.ppt) is not supported. "
                "Please save as .pptx (File → Save As → PowerPoint Presentation) and re-upload."
            )
            return doc

        try:
            from pptx import Presentation
        except ImportError:
            doc.ok = False
            doc.flag_reason = "python-pptx not installed. Run: pip install python-pptx"
            return doc

        try:
            prs = Presentation(io.BytesIO(data))
        except Exception as open_err:
            doc.ok = False
            doc.flag_reason = f"Cannot open PPTX file (may be corrupt or password-protected): {open_err}"
            return doc

        slides_data: list[dict] = []
        all_text:    list[str]  = []

        for idx, slide in enumerate(prs.slides, start=1):
            slide_text: list[str] = []
            title = ""
            for shape in slide.shapes:
                try:
                    if not shape.has_text_frame:
                        continue
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_text.append(text)
                            all_text.append(text)
                            if not title:
                                title = text
                except Exception:
                    continue  # skip shapes that crash (images, charts, etc.)

            slides_data.append({
                "slide_number": idx,
                "title":        title or f"Slide {idx}",
                "content":      " | ".join(slide_text) if slide_text else "(no text)",
            })

        if not slides_data:
            doc.ok = False
            doc.flag_reason = "PPTX file has no slides or all slides contain only images."
            return doc

        doc.raw_text   = "\n".join(all_text)
        doc.page_count = len(prs.slides)
        df             = pd.DataFrame(slides_data)
        doc.df         = df
        doc.row_count  = len(df)
        doc.col_count  = len(df.columns)
        doc.columns    = df.columns.tolist()
        doc.preview    = df.head(5).fillna("").astype(str).values.tolist()
        doc.confidence  = 0.82
        print(f"[DocumentExtractor] PPTX: {len(prs.slides)} slides")
        return doc

    # ── TXT / MD / RTF ────────────────────────────────────────────────────────
    def _parse_txt(self, data: bytes, doc: StandardDocument) -> StandardDocument:
        text = None
        for enc in ["utf-8-sig", "utf-8", "latin-1", "cp1252", "utf-16"]:
            try:
                text = data.decode(enc)
                # BOM cleanup
                text = text.lstrip("\ufeff")
                break
            except (UnicodeDecodeError, LookupError):
                continue

        if text is None:
            doc.ok = False
            doc.flag_reason = "Cannot decode text file. Try saving as UTF-8."
            return doc

        if not text.strip():
            doc.ok = False
            doc.flag_reason = "Text file is empty."
            return doc

        doc.raw_text = text
        lines = [l.strip() for l in text.split("\n") if l.strip()]
        df    = pd.DataFrame({"line_number": range(1, len(lines) + 1), "content": lines})
        doc.df         = df
        doc.row_count  = len(df)
        doc.col_count  = 2
        doc.columns    = ["line_number", "content"]
        doc.preview    = df.head(5).fillna("").astype(str).values.tolist()
        doc.confidence  = 0.85
        return doc
